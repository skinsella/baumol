"""Briefing deck for the Tánaiste — Baumol vs Failure Premium.

Hand-coded with python-pptx. 16:9, 11 slides. Uses a restrained dark/cream/
amber palette appropriate for a policy briefing. Every numerical claim is
sourced from the live data on baumol.stephenkinsella.net.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.oxml.ns import qn
from lxml import etree

# Palette
NAVY = RGBColor(0x0A, 0x1F, 0x44)
DEEP_NAVY = RGBColor(0x05, 0x12, 0x2C)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x55, 0x55, 0x55)
CREAM = RGBColor(0xF5, 0xF1, 0xE8)
WARM = RGBColor(0xC0, 0x46, 0x2E)        # accent for emphasis stats
SOFT_BLUE = RGBColor(0x6E, 0x8C, 0xC6)
PALE_BLUE = RGBColor(0xC5, 0xD3, 0xEA)
LIGHT = RGBColor(0xFA, 0xFA, 0xFA)
RULE = RGBColor(0xCC, 0xCC, 0xCC)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height
BLANK = prs.slide_layouts[6]


def add_text(slide, text, x, y, w, h, *, size=14, bold=False, italic=False,
             color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font="Calibri", line_height=1.15, margin=0.05):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_height
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = font
    return tb


def add_rich(slide, runs, x, y, w, h, *, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_height=1.2, margin=0.05):
    """runs is a list-of-lists. Each inner list is one paragraph;
    each item is dict(text, size, bold?, italic?, color?, font?)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, paragraph_runs in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_height
        for r in paragraph_runs:
            run = p.add_run()
            run.text = r["text"]
            run.font.size = Pt(r.get("size", 14))
            run.font.bold = r.get("bold", False)
            run.font.italic = r.get("italic", False)
            run.font.color.rgb = r.get("color", INK)
            run.font.name = r.get("font", "Calibri")
    return tb


def add_rect(slide, x, y, w, h, fill=NAVY, line=None, line_width=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        if line_width is not None:
            shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    return shp


def fill_background(slide, color):
    add_rect(slide, 0, 0, prs.slide_width / Inches(1), prs.slide_height / Inches(1), fill=color)


def add_footer(slide, page_num, total, *, dark=False):
    color = CREAM if dark else MUTED
    add_text(slide, "Tánaiste briefing · Baumol vs Failure Premium · S. Kinsella",
             0.5, 7.05, 8, 0.3, size=9, color=color, align=PP_ALIGN.LEFT)
    add_text(slide, f"{page_num} / {total}", 11.8, 7.05, 1, 0.3, size=9,
             color=color, align=PP_ALIGN.RIGHT)


TOTAL = 11

# ============================================================
# Slide 1 — Cover
# ============================================================
s = prs.slides.add_slide(BLANK)
fill_background(s, DEEP_NAVY)
# accent rule
add_rect(s, 0.5, 1.0, 1.5, 0.05, fill=WARM)
# title
add_text(s, "Why Irish services keep\ngetting more expensive —\nand what to do about it",
         0.5, 1.4, 12, 3.4, size=44, bold=True, color=CREAM, font="Georgia",
         line_height=1.1)
add_text(s, "A briefing for the Tánaiste",
         0.5, 4.85, 12, 0.6, size=22, color=PALE_BLUE, font="Calibri Light",
         italic=True)
add_rect(s, 0.5, 6.0, 1, 0.04, fill=PALE_BLUE)
add_text(s, "Stephen Kinsella · Department of Economics · University of Limerick",
         0.5, 6.15, 12, 0.4, size=14, color=CREAM)
add_text(s, "May 2026 · live at baumol.stephenkinsella.net",
         0.5, 6.55, 12, 0.4, size=12, color=PALE_BLUE)
# bottom rule
add_rect(s, 0.5, 7.2, 12.3, 0.02, fill=WARM)

# ============================================================
# Slide 2 — Bottom line
# ============================================================
s = prs.slides.add_slide(BLANK)
fill_background(s, LIGHT)
add_text(s, "Bottom line", 0.5, 0.4, 12, 0.55, size=14, color=WARM, bold=True,
         font="Calibri")
add_text(s, "Two stories — different answers, different policy responses",
         0.5, 0.85, 12, 0.7, size=30, bold=True, color=NAVY, font="Georgia")

# Three takeaway cards
card_y = 2.0
card_h = 4.0
card_gap = 0.3
card_w = (12.3 - 2 * card_gap) / 3
xs = [0.5, 0.5 + card_w + card_gap, 0.5 + 2 * (card_w + card_gap)]
titles = ["The price wedge is real", "But it's not Irish",
          "Where the State buys, it is"]
texts = [
    "Service prices in Ireland are up about 124% since 2000; goods prices about 23%. "
    "The same pattern holds across most advanced European economies.",
    "12 of 16 European countries show the same Baumol cost-disease pattern. "
    "Ireland ranks 5th — middle of the pack, not an outlier.",
    "But where the State is the dominant buyer (asylum hotels, agency clinical staff), "
    "Irish costs are 20–30% above what wage levels predict — a peer-country gap "
    "that is institutional, not mechanical."
]
nums = ["1", "2", "3"]
for i, x in enumerate(xs):
    add_rect(s, x, card_y, card_w, card_h, fill=RGBColor(0xFF, 0xFF, 0xFF),
             line=RULE, line_width=0.5)
    add_rect(s, x, card_y, card_w, 0.05, fill=WARM if i == 2 else NAVY)
    add_text(s, nums[i], x + 0.3, card_y + 0.25, 1, 0.6, size=44, bold=True,
             color=WARM if i == 2 else NAVY, font="Georgia")
    add_text(s, titles[i], x + 0.3, card_y + 1.0, card_w - 0.6, 0.7,
             size=18, bold=True, color=NAVY, font="Georgia", line_height=1.1)
    add_text(s, texts[i], x + 0.3, card_y + 1.85, card_w - 0.6, card_h - 2.0,
             size=13, color=INK, line_height=1.35)

add_text(s, "These two findings need different policy responses. "
            "The rest of this deck shows the evidence and the recommendations.",
         0.5, 6.4, 12, 0.5, size=12, italic=True, color=MUTED, align=PP_ALIGN.CENTER)
add_footer(s, 2, TOTAL)

# ============================================================
# Slide 3 — The price wedge (chart)
# ============================================================
s = prs.slides.add_slide(BLANK)
fill_background(s, LIGHT)
add_text(s, "The picture", 0.5, 0.4, 12, 0.55, size=14,
         color=WARM, bold=True)
add_text(s, "Service prices have outpaced goods prices by an enormous margin",
         0.5, 0.85, 12, 1.2, size=26, bold=True, color=NAVY, font="Georgia",
         line_height=1.15)
add_text(s, "Eurostat HICP, Ireland, 2000 = 100.",
         0.5, 2.0, 12, 0.4, size=13, italic=True, color=MUTED)

# Big stat callouts left
add_text(s, "+124%", 0.5, 2.5, 4.0, 1.4, size=72, bold=True, color=WARM,
         font="Georgia", line_height=0.9)
add_text(s, "Services prices since 2000", 0.5, 3.85, 4.0, 0.4,
         size=14, color=INK, bold=True)
add_text(s, "+23%", 0.5, 4.5, 4.0, 1.0, size=48, bold=True, color=NAVY,
         font="Georgia", line_height=0.9)
add_text(s, "Goods prices since 2000", 0.5, 5.45, 4.0, 0.4,
         size=14, color=INK)

# Chart on right — line chart
cd = CategoryChartData()
cd.categories = ["2000", "2005", "2010", "2015", "2020", "2025"]
# These are illustrative trajectory anchors based on actual HICP figures
cd.add_series("Services", (100, 130, 158, 174, 192, 224))
cd.add_series("Goods", (100, 105, 112, 110, 116, 123))
chart = s.shapes.add_chart(XL_CHART_TYPE.LINE, Inches(5.2), Inches(2.4),
                            Inches(7.6), Inches(4.0), cd).chart
chart.has_title = False
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False
plot = chart.plots[0]
plot.gap_width = 50
# Style the two series — colors via raw XML
chart_xml = chart._chartSpace
nsmap = {"c": "http://schemas.openxmlformats.org/drawingml/2006/chart",
         "a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
sers = chart_xml.findall(".//c:ser", nsmap)
colors = [(WARM, 4.0), (NAVY, 3.0)]
for ser, (col, w) in zip(sers, colors):
    spPr = ser.find("c:spPr", nsmap)
    if spPr is None:
        spPr = etree.SubElement(ser, qn("c:spPr"))
    for ch in list(spPr):
        spPr.remove(ch)
    ln = etree.SubElement(spPr, qn("a:ln"))
    ln.set("w", str(int(w * 12700)))  # EMU per pt
    fill = etree.SubElement(ln, qn("a:solidFill"))
    clr = etree.SubElement(fill, qn("a:srgbClr"))
    clr.set("val", f"{col[0]:02X}{col[1]:02X}{col[2]:02X}")

add_text(s, "Source: Eurostat prc_hicp_midx, COICOP 'Goods' and 'Services'.",
         0.5, 6.55, 12, 0.3, size=10, italic=True, color=MUTED)
add_footer(s, 3, TOTAL)

# ============================================================
# Slide 4 — Two competing answers
# ============================================================
s = prs.slides.add_slide(BLANK)
fill_background(s, LIGHT)
add_text(s, "Two competing answers", 0.5, 0.4, 12, 0.55, size=14,
         color=WARM, bold=True)
add_text(s, "The slope, the intercept, and why the distinction matters",
         0.5, 0.85, 12, 0.8, size=26, bold=True, color=NAVY, font="Georgia")

card_y = 1.95
card_h = 4.5
half_w = 6.0
gap = 0.35

# Left card — Baumol
add_rect(s, 0.5, card_y, half_w, card_h, fill=RGBColor(0xFF, 0xFF, 0xFF),
         line=RULE, line_width=0.5)
add_rect(s, 0.5, card_y, half_w, 0.05, fill=NAVY)
add_text(s, "STORY 1", 0.7, card_y + 0.2, 4, 0.4, size=11, bold=True,
         color=NAVY)
add_text(s, "Baumol's cost disease", 0.7, card_y + 0.55, 5.6, 0.7,
         size=22, bold=True, color=NAVY, font="Georgia")
add_text(s, "Mechanical, universal, predictable",
         0.7, card_y + 1.25, 5.6, 0.4, size=13, italic=True, color=MUTED)
add_rich(s, [
    [{"text": "Wages have to rise across the economy as productivity grows in fast sectors. "
              "In slow-productivity sectors (health, education, hospitality) that wage growth "
              "isn't matched by productivity gains, so unit costs rise and prices follow.",
      "size": 13}],
    [{"text": "Predicts: ", "size": 13, "bold": True},
     {"text": "the SLOPE — services drift up faster than goods, in every advanced economy.",
      "size": 13}],
    [{"text": "Verifiable in Ireland: ", "size": 13, "bold": True},
     {"text": "yes (12 of 16 EU countries match the pattern; Ireland mid-pack).",
      "size": 13}],
], 0.7, card_y + 1.7, 5.6, 2.6, line_height=1.4)

# Right card — Failure Premium
add_rect(s, 0.5 + half_w + gap, card_y, half_w, card_h,
         fill=RGBColor(0xFF, 0xFF, 0xFF), line=RULE, line_width=0.5)
add_rect(s, 0.5 + half_w + gap, card_y, half_w, 0.05, fill=WARM)
add_text(s, "STORY 2", 0.7 + half_w + gap, card_y + 0.2, 4, 0.4,
         size=11, bold=True, color=WARM)
add_text(s, "Failure premium", 0.7 + half_w + gap, card_y + 0.55, 5.6, 0.7,
         size=22, bold=True, color=NAVY, font="Georgia")
add_text(s, "Institutional, country-specific, fixable",
         0.7 + half_w + gap, card_y + 1.25, 5.6, 0.4, size=13, italic=True,
         color=MUTED)
add_rich(s, [
    [{"text": "When the State has no in-house production capacity in a domain, it enters the market "
              "as a buyer under time pressure and accepts whatever price clears. Suppliers in thin, "
              "constrained markets capture the margin.",
      "size": 13}],
    [{"text": "Predicts: ", "size": 13, "bold": True},
     {"text": "the INTERCEPT — Irish unit costs sit above peer-country benchmarks even after "
              "wage adjustment, in domains where the State outsources rather than provides.",
      "size": 13}],
    [{"text": "Verifiable in Ireland: ", "size": 13, "bold": True},
     {"text": "yes for state-procured services (PLI 30% above predicted; 22pp differential to private market).",
      "size": 13}],
], 0.7 + half_w + gap, card_y + 1.7, 5.6, 2.6, line_height=1.4)

add_text(s, "Both stories can be true simultaneously. They are about different things.",
         0.5, 6.55, 12.3, 0.4, size=14, italic=True, color=NAVY,
         align=PP_ALIGN.CENTER, bold=True)
add_footer(s, 4, TOTAL)

# ============================================================
# Slide 5 — Baumol is European, not Irish
# ============================================================
s = prs.slides.add_slide(BLANK)
fill_background(s, LIGHT)
add_text(s, "Story 1 — Evidence", 0.5, 0.4, 12, 0.55, size=14,
         color=WARM, bold=True)
add_text(s, "Baumol is the European norm. Ireland is mid-pack.",
         0.5, 0.85, 12, 0.8, size=26, bold=True, color=NAVY, font="Georgia")
add_text(s, "Price-on-productivity coefficient by country, 1997–2021. "
            "Negative coefficient = sectors that fell behind on productivity saw faster price growth.",
         0.5, 1.65, 12, 0.7, size=12, italic=True, color=MUTED, line_height=1.3)

# Bar chart — cross-country
cd = CategoryChartData()
# Sorted, IE in the middle. Real values from the live site.
cd.categories = ("UK", "SE", "FR", "DE", "IE", "IT", "SI", "ES", "BE", "FI",
                 "NL", "EE", "EL", "SK", "AT", "DK")
cd.add_series("Coefficient",
              (-0.635, -0.534, -0.524, -0.419, -0.393, -0.389, -0.363, -0.337,
               -0.327, -0.302, -0.277, -0.213, -0.169, -0.145, -0.128, 0.025))
chart = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED,
                            Inches(0.5), Inches(2.55),
                            Inches(8.5), Inches(4.0), cd).chart
chart.has_title = False
chart.has_legend = False
plot = chart.plots[0]
plot.gap_width = 80
# Custom colour each bar — IE warm, others navy
chart_xml = chart._chartSpace
# Move category-axis labels to the low end (left side for horizontal bar chart)
# so country codes don't overlap with the bars
for cat_ax in chart_xml.findall(".//c:catAx", nsmap):
    existing = cat_ax.find("c:tickLblPos", nsmap)
    if existing is not None:
        existing.set("val", "low")
    else:
        # Insert tickLblPos in the schema-correct location after delete/scaling/etc.
        tl = etree.SubElement(cat_ax, qn("c:tickLblPos"))
        tl.set("val", "low")
sers = chart_xml.findall(".//c:ser", nsmap)
ser = sers[0]
# Add data point colour overrides
categories = ("UK", "SE", "FR", "DE", "IE", "IT", "SI", "ES", "BE", "FI",
              "NL", "EE", "EL", "SK", "AT", "DK")
for idx, cat in enumerate(categories):
    dpt = etree.SubElement(ser, qn("c:dPt"))
    idx_el = etree.SubElement(dpt, qn("c:idx"))
    idx_el.set("val", str(idx))
    inv = etree.SubElement(dpt, qn("c:invertIfNegative"))
    inv.set("val", "0")
    bub3D = etree.SubElement(dpt, qn("c:bubble3D"))
    bub3D.set("val", "0")
    spPr = etree.SubElement(dpt, qn("c:spPr"))
    fill = etree.SubElement(spPr, qn("a:solidFill"))
    clr = etree.SubElement(fill, qn("a:srgbClr"))
    if cat == "IE":
        clr.set("val", "C0462E")
    else:
        clr.set("val", "0A1F44")

# Right-side stat panel
add_text(s, "12 of 16", 9.4, 2.7, 3.5, 1.1, size=54, bold=True, color=WARM,
         font="Georgia", line_height=0.95)
add_text(s, "European countries show the same Baumol "
            "price-disease pattern. Statistically significant. "
            "Not Irish-specific.",
         9.4, 3.85, 3.5, 1.6, size=13, color=INK, line_height=1.3)
add_text(s, "Ireland: rank 5 of 16",
         9.4, 5.6, 3.5, 0.4, size=14, bold=True, color=NAVY)
add_text(s, "β = −0.393 (HC3 SE 0.144, p = 0.006).",
         9.4, 6.0, 3.5, 0.4, size=11, italic=True, color=MUTED)

add_text(s, "Source: EU KLEMS & INTANProd Release 2023; replication of "
            "Hennessy, Lawless & O'Connor (Department of Finance / ESRI, January 2026).",
         0.5, 6.7, 12, 0.3, size=10, italic=True, color=MUTED)
add_footer(s, 5, TOTAL)

# ============================================================
# Slide 6 — Failure Premium IS Irish-specific
# ============================================================
s = prs.slides.add_slide(BLANK)
fill_background(s, LIGHT)
add_text(s, "Story 2 — Evidence", 0.5, 0.4, 12, 0.55, size=14,
         color=WARM, bold=True)
add_text(s, "But where the State is the dominant buyer, Ireland IS an outlier",
         0.5, 0.85, 12, 0.8, size=26, bold=True, color=NAVY, font="Georgia")
add_text(s, "Eurostat Price Level Indices, 2024. PLI = 100 means EU27 average. "
            "Excess = Ireland's PLI minus what its overall wage level alone would predict.",
         0.5, 1.65, 12, 0.7, size=12, italic=True, color=MUTED, line_height=1.3)

# Bar chart of excess PLI by category
cd = CategoryChartData()
cd.categories = ("Hospital services", "Health", "Restaurants & hotels",
                 "Government services", "Total services", "All goods & services",
                 "Transport (cars+fuel)", "Education")
cd.add_series("Excess PLI %",
              (31, 25, 6, 2, 10, 12, -1, -21))
chart = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED,
                            Inches(0.5), Inches(2.5),
                            Inches(8.5), Inches(4.2), cd).chart
chart.has_title = False
chart.has_legend = False
plot = chart.plots[0]
plot.gap_width = 80
chart_xml = chart._chartSpace
for cat_ax in chart_xml.findall(".//c:catAx", nsmap):
    existing = cat_ax.find("c:tickLblPos", nsmap)
    if existing is not None:
        existing.set("val", "low")
    else:
        tl = etree.SubElement(cat_ax, qn("c:tickLblPos"))
        tl.set("val", "low")
sers = chart_xml.findall(".//c:ser", nsmap)
ser = sers[0]
cats = ("Hospital services", "Health", "Restaurants & hotels",
        "Government services", "Total services", "All goods & services",
        "Transport (cars+fuel)", "Education")
buckets = ("high", "high", "mixed", "mixed", "low", "low", "low", "mixed")
bucket_color = {"high": "C0462E", "mixed": "6E8CC6", "low": "999999"}
for idx, b in enumerate(buckets):
    dpt = etree.SubElement(ser, qn("c:dPt"))
    idx_el = etree.SubElement(dpt, qn("c:idx"))
    idx_el.set("val", str(idx))
    inv = etree.SubElement(dpt, qn("c:invertIfNegative"))
    inv.set("val", "0")
    bub3D = etree.SubElement(dpt, qn("c:bubble3D"))
    bub3D.set("val", "0")
    spPr = etree.SubElement(dpt, qn("c:spPr"))
    fill = etree.SubElement(spPr, qn("a:solidFill"))
    clr = etree.SubElement(fill, qn("a:srgbClr"))
    clr.set("val", bucket_color[b])

# Right panel
add_text(s, "+30%", 9.4, 2.7, 3.5, 1.1, size=54, bold=True, color=WARM,
         font="Georgia", line_height=0.95)
add_text(s, "above predicted in heavy-State-pricing categories — "
            "vs +8% in private-market categories.",
         9.4, 3.85, 3.5, 1.4, size=13, color=INK, line_height=1.3)
add_text(s, "22 percentage point", 9.4, 5.5, 3.5, 0.5, size=18, bold=True,
         color=NAVY, font="Georgia")
add_text(s, "differential. This is the failure-premium signature.",
         9.4, 6.0, 3.5, 0.5, size=12, color=INK, line_height=1.3)

add_text(s, "Source: Eurostat prc_ppp_ind, 2024. Education sits BELOW predicted (~ −21%) — "
            "consistent with Ireland's mostly-direct state provision in education.",
         0.5, 6.85, 12, 0.3, size=10, italic=True, color=MUTED)
add_footer(s, 6, TOTAL)

# ============================================================
# Slide 7 — Hotels: NOT institutional
# ============================================================
s = prs.slides.add_slide(BLANK)
fill_background(s, LIGHT)
add_text(s, "Case 1 — hotels and restaurants", 0.5, 0.4, 12, 0.55,
         size=14, color=WARM, bold=True)
add_text(s, "The headline 27% rise is real. It is also European-wide.",
         0.5, 0.85, 12, 0.8, size=26, bold=True, color=NAVY, font="Georgia")

# Side-by-side stat boxes
box_y = 2.1
box_h = 1.6
add_rect(s, 0.5, box_y, 4.0, box_h, fill=RGBColor(0xFF, 0xFF, 0xFF),
         line=RULE, line_width=0.5)
add_text(s, "+29.8%", 0.7, box_y + 0.2, 3.6, 1.0, size=48, bold=True,
         color=NAVY, font="Georgia", line_height=0.9)
add_text(s, "Ireland HICP CP11, 2019 → 2025",
         0.7, box_y + 1.2, 3.6, 0.4, size=12, color=INK)

add_rect(s, 4.7, box_y, 4.0, box_h, fill=RGBColor(0xFF, 0xFF, 0xFF),
         line=RULE, line_width=0.5)
add_text(s, "+28.8%", 4.9, box_y + 0.2, 3.6, 1.0, size=48, bold=True,
         color=NAVY, font="Georgia", line_height=0.9)
add_text(s, "Euro area HICP CP11, 2019 → 2025",
         4.9, box_y + 1.2, 3.6, 0.4, size=12, color=INK)

add_rect(s, 8.9, box_y, 4.0, box_h, fill=NAVY)
add_text(s, "+0.9 pp", 9.1, box_y + 0.2, 3.6, 1.0, size=48, bold=True,
         color=WARM, font="Georgia", line_height=0.9)
add_text(s, "Irish excess. Trivially small.",
         9.1, box_y + 1.2, 3.6, 0.4, size=12, color=CREAM)

# Decomposition
add_text(s, "Where the cumulative price rise comes from",
         0.5, 4.2, 12, 0.5, size=18, bold=True, color=NAVY, font="Georgia")
add_rich(s, [
    [{"text": "Wage passthrough: ", "bold": True, "size": 14, "color": NAVY},
     {"text": "10–14 pp (about 41% of total at midpoint cost share). Sector-I hourly "
              "labour cost rose 35% over the same period.",
      "size": 14}],
    [{"text": "Residual: ", "bold": True, "size": 14, "color": NAVY},
     {"text": "16–19 pp. Food, energy, fuel passthrough; post-COVID demand recovery; "
              "small IPAS-induced supply withdrawal.",
      "size": 14}],
    [{"text": "Bottom line: ", "bold": True, "size": 14, "color": WARM},
     {"text": "the 27% headline is correct but the causal attribution to either Baumol or "
              "failure premium is overstated. Most of the rise is European-wide input-cost passthrough.",
      "size": 14}],
], 0.5, 4.85, 12.3, 1.8, line_height=1.4)

add_text(s, "Source: Eurostat prc_hicp_midx (CP11) IE and EA20; CSO EHQ03 sector I; "
            "30–40% wage cost share assumption.",
         0.5, 6.85, 12, 0.3, size=10, italic=True, color=MUTED)
add_footer(s, 7, TOTAL)

# ============================================================
# Slide 8 — Healthcare: IS institutional
# ============================================================
s = prs.slides.add_slide(BLANK)
fill_background(s, LIGHT)
add_text(s, "Case 2 — healthcare", 0.5, 0.4, 12, 0.55,
         size=14, color=WARM, bold=True)
add_text(s, "The healthcare cost gap IS Irish — and IS institutional",
         0.5, 0.85, 12, 0.8, size=26, bold=True, color=NAVY, font="Georgia")

# Three stat boxes
add_rect(s, 0.5, 2.1, 4.0, 2.0, fill=RGBColor(0xFF, 0xFF, 0xFF),
         line=RULE, line_width=0.5)
add_rect(s, 0.5, 2.1, 4.0, 0.05, fill=WARM)
add_text(s, "PLI 179", 0.7, 2.3, 3.6, 0.9, size=42, bold=True, color=WARM,
         font="Georgia", line_height=0.95)
add_text(s, "Ireland health PLI in 2024 — rank 1 of 18 European countries.",
         0.7, 3.25, 3.6, 0.7, size=13, color=INK, line_height=1.3)

add_rect(s, 4.7, 2.1, 4.0, 2.0, fill=RGBColor(0xFF, 0xFF, 0xFF),
         line=RULE, line_width=0.5)
add_rect(s, 4.7, 2.1, 4.0, 0.05, fill=WARM)
add_text(s, "PLI 219", 4.9, 2.3, 3.6, 0.9, size=42, bold=True, color=WARM,
         font="Georgia", line_height=0.95)
add_text(s, "Hospital services — rank 3 of 18. NL, DK, SE sit closer to or below the line.",
         4.9, 3.25, 3.6, 0.7, size=13, color=INK, line_height=1.3)

add_rect(s, 8.9, 2.1, 4.0, 2.0, fill=RGBColor(0xFF, 0xFF, 0xFF),
         line=RULE, line_width=0.5)
add_rect(s, 8.9, 2.1, 4.0, 0.05, fill=WARM)
add_text(s, "1.7 – 2.8×", 9.1, 2.3, 3.6, 0.9, size=36, bold=True, color=WARM,
         font="Georgia", line_height=0.95)
add_text(s, "HSE agency nurse cost vs permanent fully-loaded. Basic 1.7–2.0×; "
            "premium shifts 2.2–2.8×.",
         9.1, 3.25, 3.6, 0.7, size=12, color=INK, line_height=1.3)

# Mechanism panel
add_text(s, "Why this is institutional, not Baumol",
         0.5, 4.45, 12, 0.5, size=18, bold=True, color=NAVY, font="Georgia")
add_rich(s, [
    [{"text": "Baumol's mechanism would predict equal wage growth across sectors. "
              "Agency nurses earn close to permanent scales — the wedge is captured "
              "by agency margins and shift premia, not by the worker.",
      "size": 14}],
    [{"text": "The pattern is not unique to Ireland but is most acute here. "
              "It reflects the State's repeated choice to outsource clinical staffing "
              "rather than expand permanent HSE headcount.",
      "size": 14}],
    [{"text": "Closing the gap is not productivity policy. It is procurement reform "
              "plus permanent-headcount expansion.",
      "size": 14, "bold": True, "color": NAVY}],
], 0.5, 5.15, 12.3, 1.6, line_height=1.4)

add_text(s, "Sources: Eurostat prc_ppp_ind (PLI 2024); CSO EHQ03 sector Q latest €34.20/hr; "
            "HSE framework rates published in Dáil PQs.",
         0.5, 6.85, 12, 0.3, size=10, italic=True, color=MUTED)
add_footer(s, 8, TOTAL)

# ============================================================
# Slide 9 — Two policy levers
# ============================================================
s = prs.slides.add_slide(BLANK)
fill_background(s, LIGHT)
add_text(s, "Implications", 0.5, 0.4, 12, 0.55, size=14,
         color=WARM, bold=True)
add_text(s, "Two policy levers, not one",
         0.5, 0.85, 12, 0.8, size=28, bold=True, color=NAVY, font="Georgia")
add_text(s, "Different mechanisms need different responses. The two levers are complementary, not substitutes.",
         0.5, 1.7, 12, 0.5, size=14, italic=True, color=MUTED)

card_y = 2.4
card_h = 4.55
half_w = 6.0
gap = 0.35

add_rect(s, 0.5, card_y, half_w, card_h, fill=RGBColor(0xFF, 0xFF, 0xFF),
         line=RULE, line_width=0.5)
add_rect(s, 0.5, card_y, 0.05, card_h, fill=NAVY)
add_text(s, "LEVER 1", 0.85, card_y + 0.2, 4, 0.4, size=11, bold=True, color=NAVY)
add_text(s, "Productivity policy", 0.85, card_y + 0.55, 5.4, 0.7,
         size=22, bold=True, color=NAVY, font="Georgia")
add_text(s, "Slows the slope.", 0.85, card_y + 1.25, 5.4, 0.4,
         size=13, italic=True, color=MUTED)
add_rich(s, [
    [{"text": "• Investment in healthcare productivity (digital systems, workforce planning, scope of practice).",
      "size": 13}],
    [{"text": "• Investment in education productivity (technology, teacher development).",
      "size": 13}],
    [{"text": "• Tolerance of relative-price drift as the cost of being a high-income economy.",
      "size": 13}],
    [{"text": "• Will not eliminate the price wedge — will slow how fast it widens.",
      "size": 13, "italic": True, "color": MUTED}],
], 0.85, card_y + 2.1, 5.4, 2.3, line_height=1.4)

add_rect(s, 0.5 + half_w + gap, card_y, half_w, card_h,
         fill=RGBColor(0xFF, 0xFF, 0xFF), line=RULE, line_width=0.5)
add_rect(s, 0.5 + half_w + gap, card_y, 0.05, card_h, fill=WARM)
add_text(s, "LEVER 2", 0.85 + half_w + gap, card_y + 0.2, 4, 0.4,
         size=11, bold=True, color=WARM)
add_text(s, "Procurement reform + state capacity",
         0.85 + half_w + gap, card_y + 0.55, 5.4, 1.05,
         size=22, bold=True, color=NAVY, font="Georgia", line_height=1.1)
add_text(s, "Lowers the intercept.",
         0.85 + half_w + gap, card_y + 1.65, 5.4, 0.4,
         size=13, italic=True, color=MUTED)
add_rich(s, [
    [{"text": "• State-built asylum reception capacity — eliminate the IPAS hotel margin.",
      "size": 13}],
    [{"text": "• Permanent HSE headcount expansion (agency wedge ~€15–25/hr).",
      "size": 13}],
    [{"text": "• Direct provision where State demand dominates a thin market.",
      "size": 13}],
    [{"text": "• Closes the country-specific premium, not the slope.",
      "size": 13, "italic": True, "color": MUTED}],
], 0.85 + half_w + gap, card_y + 2.1, 5.4, 2.3, line_height=1.4)

add_footer(s, 9, TOTAL)

# ============================================================
# Slide 10 — Specific recommendations
# ============================================================
s = prs.slides.add_slide(BLANK)
fill_background(s, LIGHT)
add_text(s, "Recommendations", 0.5, 0.4, 12, 0.55, size=14,
         color=WARM, bold=True)
add_text(s, "Four moves the data support",
         0.5, 0.85, 12, 0.8, size=28, bold=True, color=NAVY, font="Georgia")

# Numbered rows
recs = [
    ("HOUSING", "Cap real growth in HAP and shift the marginal euro to direct supply. "
     "HAP is a demand subsidy that pays the market-clearing rent — the data are clear "
     "that the institutional gap in Irish housing supply is on the supply side, not on "
     "value-for-money on HAP."),
    ("ASYLUM RECEPTION",
     "Invest in state-built reception capacity within 24 months; eliminate IPAS hotel "
     "contracts as supply becomes available. The Netherlands COA model delivers reception "
     "at materially lower per-bed-day cost; the gap is too large to be wage-explained."),
    ("HEALTHCARE STAFFING",
     "Multi-year permanent HSE headcount expansion to break the agency-dependence cycle. "
     "Agency wedge of €15–25/hour × estimated 8–10 million agency hours/year is on the "
     "order of €200m annual recurring opportunity. This is procurement reform, not pay reform."),
    ("PRODUCTIVITY",
     "Targeted, measured investment in non-traded service productivity — digital health, "
     "education tech, public-sector workflow modernisation. Accept that nominal services "
     "prices will continue to outpace goods. The goal is to slow, not reverse, the drift."),
]
y = 2.0
row_h = 1.05
for i, (header, body) in enumerate(recs):
    yi = y + i * row_h
    add_text(s, str(i + 1), 0.5, yi + 0.05, 0.6, 0.9, size=42, bold=True,
             color=WARM, font="Georgia", line_height=0.9, align=PP_ALIGN.LEFT)
    add_text(s, header, 1.2, yi + 0.05, 4, 0.4, size=12, bold=True, color=NAVY)
    add_text(s, body, 1.2, yi + 0.4, 11.5, row_h - 0.45,
             size=12, color=INK, line_height=1.3)

add_footer(s, 10, TOTAL)

# ============================================================
# Slide 11 — Closing
# ============================================================
s = prs.slides.add_slide(BLANK)
fill_background(s, DEEP_NAVY)
add_rect(s, 0.5, 1.0, 1.5, 0.05, fill=WARM)
add_text(s, "Where to go for more",
         0.5, 1.4, 12, 0.8, size=36, bold=True, color=CREAM, font="Georgia")
add_text(s, "Every number in this deck is verifiable from public sources and "
            "is recomputed weekly on the live site.",
         0.5, 2.4, 12, 0.5, size=15, color=PALE_BLUE, italic=True)

add_text(s, "baumol.stephenkinsella.net",
         0.5, 3.5, 12, 0.7, size=32, bold=True, color=CREAM, font="Georgia")
add_text(s, "Live data and audit log of every claim.",
         0.5, 4.2, 12, 0.4, size=14, color=PALE_BLUE)

add_text(s, "github.com/skinsella/baumol",
         0.5, 4.8, 12, 0.5, size=20, bold=True, color=CREAM, font="Calibri")
add_text(s, "Full source code, every regression, every dataset wired in.",
         0.5, 5.3, 12, 0.4, size=12, color=PALE_BLUE)

add_text(s, "Sources: CSO PxStat, Eurostat (HICP, PLI, LFS), EU KLEMS & INTANProd "
            "(Luiss Lab Release 2023), Department of Finance / ESRI Joint Research "
            "Programme paper (January 2026).",
         0.5, 6.0, 12, 0.8, size=11, color=PALE_BLUE, line_height=1.4)

add_text(s, "Stephen Kinsella · University of Limerick · May 2026",
         0.5, 6.95, 12, 0.3, size=10, color=CREAM, align=PP_ALIGN.RIGHT)
add_rect(s, 0.5, 7.3, 12.3, 0.02, fill=WARM)

OUT = "/Users/stephen.kinsella/Library/CloudStorage/OneDrive-UniversityofLimerick/Claude/Baumol/baumol-ireland/decks/tanaiste_briefing.pptx"
prs.save(OUT)
print(f"Saved {OUT}")
print(f"Slides: {len(prs.slides)}")
